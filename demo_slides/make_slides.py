#!/usr/bin/env python3
"""
Generate a PPTX slide deck for the "Disposable Web Apps with LLMs" talk.

Requires: python-pptx (pip install python-pptx)
Run from repo root:
  python demo_slides/make_slides.py
Outputs:
  demo_slides/llm_disposable_webapps_demo.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.stderr.write("Missing dependency: python-pptx\nInstall: pip install python-pptx\n")
    sys.exit(1)


TITLE = "Disposable Web Apps with LLMs"
SUBTITLE = "Tiny, task-specific UIs you can build in minutes"


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE


def add_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.add_paragraph()
        p.text = line
        p.level = 0


def add_two_column(prs, title, left, right):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # two-content
    slide.shapes.title.text = title
    tf_left = slide.placeholders[1].text_frame
    tf_left.clear()
    for item in left:
        p = tf_left.add_paragraph()
        p.text = item
        p.level = 0
    tf_right = slide.placeholders[2].text_frame
    tf_right.clear()
    for item in right:
        p = tf_right.add_paragraph()
        p.text = item
        p.level = 0


def main():
    prs = Presentation()

    add_title_slide(prs)

    add_bullets(
        prs,
        "Why this matters now",
        [
            "LLMs make single-use web apps in minutes, not days.",
            "Goal: unlock tasks you’d never do because they were too tedious.",
            "Disposable, data-aware interfaces on demand.",
        ],
    )

    add_bullets(
        prs,
        "The big unlock: disposable web apps",
        [
            "Tiny, task-specific pages that wrap your data.",
            "Faster insight: filter, slice, click, and read—no CSV spelunking.",
            "Built when needed; tossed when done; zero product overhead.",
        ],
    )

    # Example 1 slides
    add_bullets(
        prs,
        "Example 1: Labeler (UI highlights)",
        [
            "Chat bubbles (user right, AI left); per-thread meta.",
            "Inline labeling, flags, notes; save/resume; export CSV/JSON.",
            "Context window chart to stay anchored to day/time.",
        ],
    )
    add_bullets(
        prs,
        "Example 1: What it enabled",
        [
            "Read chats like chats, not rows — huge readability win.",
            "Label quickly; flag edge cases; export disagreements instantly.",
            "Filter by region/date/time window without leaving the view.",
        ],
    )
    add_bullets(
        prs,
        "Example 1: Patterns you can reuse",
        [
            "Treat rows as conversations (pair + align bubbles).",
            "Inline labeling UI + autosave + export.",
            "Mini-analytics pane beside the work surface.",
        ],
    )

    # Example 2 slides
    add_bullets(
        prs,
        "Example 2: Time-Window Explorer (UI highlights)",
        [
            "Hour filters, bars, scatter, region buckets, thread browser.",
            "Inside vs outside window shading; quick bucket navigation.",
            "Exports the filtered slice for follow-up analysis.",
        ],
    )
    add_bullets(
        prs,
        "Example 2: What it enabled",
        [
            "See which hours/regions light up at a glance.",
            "Click buckets → jump straight to threads and read them.",
            "Keep filters, charts, and detail pane in sync.",
        ],
    )
    add_bullets(
        prs,
        "Example 2: Patterns you can reuse",
        [
            "One control panel drives multiple coordinated charts.",
            "Bucketed lists + detail pane for fast drill-down.",
            "Inside/outside band highlighting to answer “when” fast.",
        ],
    )

    add_bullets(
        prs,
        "Example 3 (placeholder)",
        [
            "Another single-use tool from a different project.",
            "Show UI → list the work it makes faster/better.",
            "Extract reusable patterns attendees can copy.",
        ],
    )

    add_bullets(
        prs,
        "Prompting techniques: data briefing",
        [
            "Describe schema: columns, types, what a row represents, units.",
            "Provide 2–5 sample rows to de-risk parsing.",
            "Note PII/sensitivity rules and what NOT to store/display.",
        ],
    )

    add_bullets(
        prs,
        "Prompting techniques: project briefing",
        [
            "State objective, audience, definition of done.",
            "List constraints: time, security, offline/online, allowed libs.",
            "Clarify what you won’t do (e.g., no DB writes, no auth).",
        ],
    )

    add_bullets(
        prs,
        "Prompting techniques: tasks + surprise me",
        [
            "Give explicit asks (e.g., bar: x=hour, y=thread count).",
            "List research questions the app must answer.",
            "End with: “Surprise me with one extra view that helps answer these.”",
        ],
    )

    add_bullets(
        prs,
        "Where to get the demos",
        [
            "demo_labeler/ — chat-style manual labeler with synthetic data.",
            "demo_time_window_viewer/ — time-window explorer with synthetic data.",
            "Run: python3 -m http.server 8000 (inside the folder) → open http://localhost:8000/",
            "GitHub/internal link: <add your link>",
        ],
    )

    out_path = Path(__file__).parent / "llm_disposable_webapps_demo.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
